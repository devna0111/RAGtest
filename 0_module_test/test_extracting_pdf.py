import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import StringIO


def pptx_to_markdown_string(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    output = StringIO()
    output.write("# PPT 자동 변환\n\n")

    def recurse_shapes(shapes, output_lines):
        for shape in shapes:
            # 텍스트 추출
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    output_lines.append(f"{text}\n")

            # 이미지 추출 설명
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                output_lines.append("[이미지 포함됨]\n")

            # 그룹 도형 내부 재귀
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                recurse_shapes(shape.shapes, output_lines)

    for i, slide in enumerate(prs.slides):
        output.write(f"## 슬라이드 {i + 1}\n\n")
        slide_lines = []
        recurse_shapes(slide.shapes, slide_lines)
        output.write("".join(slide_lines))
        output.write("\n---\n\n")

    return output.getvalue()


if __name__ == "__main__":
    FILE_PATH = "input_slides.pptx"
    # OUTPUT_PATH = "output_from_ppt.md"

    result = pptx_to_markdown_string(FILE_PATH)
    print(result)

    # with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
    #     f.write(result)

    # print(f"[저장 완료] {OUTPUT_PATH}")
