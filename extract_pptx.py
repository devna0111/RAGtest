import pptx
import os
from PIL import Image
import io
from paddleocr import PaddleOCR

ocr = PaddleOCR(lang="korean")

def extract_pptx_chunks_with_images(filepath, img_dir='img_output_pptx'):
    prs = pptx.Presentation(filepath)
    os.makedirs(img_dir, exist_ok=True)
    chunks = []
    image_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # 텍스트 추출
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    chunks.append({"type": "text", "content": text})

            # 표 추출
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                table = shape.table
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append('\t'.join(row_text))
                tstr = "\n".join(table_text)
                if tstr.strip():
                    chunks.append({"type": "table", "content": tstr})

            # 이미지 추출
            if 'Picture' in str(type(shape)):
                img = shape.image
                image_bytes = img.blob
                img_path = os.path.join(img_dir, f"slide{slide_idx}_img{image_count}.png")
                image_count += 1
                with open(img_path, "wb") as f:
                    f.write(image_bytes)
                # PIL로 강제 변환/재저장(에러 방지)
                try:
                    pil_img = Image.open(img_path)
                    pil_img = pil_img.convert("RGB")
                    pil_img.save(img_path)
                    chunks.append({"type": "image", "content": img_path})
                    # OCR
                    try:
                        result = ocr.predict(img_path)
                        ocr_lines = result[0]['rec_texts'] if result else []
                        ocr_text = "\n".join(ocr_lines)
                        if ocr_text:
                            chunks.append({"type": "ocr", "content": ocr_text})
                    except Exception as ocr_err:
                        print(f"OCR 실패: {img_path} - {ocr_err}")
                except Exception as img_err:
                    print(f"이미지 변환 실패: {img_path} - {img_err}")

    return chunks

if __name__ == "__main__":
    out = extract_pptx_chunks_with_images("sample.pptx")
    for chunk in out:
        print(f"[{chunk['type']}]", chunk['content'][:80])
