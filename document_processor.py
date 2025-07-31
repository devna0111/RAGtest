import os
import base64
import fitz  # PyMuPDF
from PIL import Image
import io

from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

from langchain_community.llms import Ollama
from langchain.schema.messages import HumanMessage

# --- 설정 ---
VISION_MODEL = "qwen2.5vl:7b"
TEMP_IMAGE_DIR = "temp_images"
os.makedirs(TEMP_IMAGE_DIR, exist_ok=True)

# --- Vision 모델 호출 함수 ---
def get_visual_description(image_bytes: bytes, prompt: str) -> str:
    """
    이미지 바이트 데이터를 받아 Vision LLM으로 분석하고 설명을 반환합니다.
    """
    try:
        llm = Ollama(model=VISION_MODEL)
        
        # 이미지를 base64로 인코딩
        b64_image = base64.b64encode(image_bytes).decode('utf-8')
        
        # HumanMessage 구성
        msg = llm.invoke(
            [
                HumanMessage(
                    content=[
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": f"data:image/jpeg;base64,{b64_image}",
                        },
                    ]
                )
            ]
        )
        return msg
    except Exception as e:
        print(f"Vision 모델 호출 중 오류 발생: {e}")
        return "이미지 분석 중 오류가 발생했습니다."

# --- 파일 유형별 처리 함수 ---
def _process_pdf(filepath: str) -> list[str]:
    """PDF에서 텍스트와 이미지를 추출합니다."""
    texts = []
    doc = fitz.open(filepath)
    for page_num, page in enumerate(doc):
        # 텍스트 추출
        texts.append(page.get_text())
        
        # 이미지 추출
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            
            prompt = f"이것은 문서 {page_num+1}페이지에 있는 이미지입니다. 이 이미지(표 또는 그래프일 수 있음)를 자세히 설명해주세요. 마크다운 형식으로 답변해주세요."
            desc = get_visual_description(image_bytes, prompt)
            texts.append(f"\n[페이지 {page_num+1}의 이미지 설명]:\n{desc}\n")
    doc.close()
    return texts

def _process_docx(filepath: str) -> list[str]:
    """DOCX에서 텍스트와 이미지를 추출합니다."""
    texts = []
    doc = DocxDocument(filepath)
    
    # 텍스트 추출
    for para in doc.paragraphs:
        texts.append(para.text)
        
    # 이미지 추출 (python-docx는 직접적인 이미지 추출이 복잡하므로, 이 부분은 고도화 과제로 남깁니다)
    # 현재는 텍스트만 추출합니다.
    texts.append("\n[알림] 현재 버전에서는 DOCX 파일의 텍스트만 추출합니다.\n")
    
    return texts

def _process_pptx(filepath: str) -> list[str]:
    """PPTX에서 텍스트와 이미지를 추출합니다."""
    texts = []
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        slide_texts = [f"\n--- 슬라이드 {i+1} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_texts.append(shape.text)
            if shape.shape_type == 13: # 이미지(Picture)
                image_bytes = shape.image.blob
                prompt = f"이것은 프레젠테이션 {i+1}번째 슬라이드의 이미지입니다. 이 이미지(표 또는 그래프일 수 있음)를 자세히 설명해주세요. 마크다운 형식으로 답변해주세요."
                desc = get_visual_description(image_bytes, prompt)
                slide_texts.append(f"\n[슬라이드 {i+1}의 이미지 설명]:\n{desc}\n")
        texts.extend(slide_texts)
    return texts

def _process_xlsx(filepath: str) -> list[str]:
    """XLSX에서 셀 텍스트를 추출합니다."""
    texts = []
    workbook = load_workbook(filepath)
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        sheet_text = f"\n--- 시트: {sheet_name} ---\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) for cell in row if cell is not None])
            if row_text:
                sheet_text += row_text + "\n"
        texts.append(sheet_text)
    texts.append("\n[알림] 현재 버전에서는 XLSX 파일의 텍스트 데이터만 추출합니다. 차트/그래프 분석은 지원하지 않습니다.\n")
    return texts

def _process_image(filepath: str) -> list[str]:
    """단일 이미지 파일을 처리합니다."""
    with open(filepath, "rb") as f:
        image_bytes = f.read()
    prompt = "이 이미지를 자세히 설명해주세요. 마크다운 형식으로 답변해주세요."
    desc = get_visual_description(image_bytes, prompt)
    return [f"[이미지 파일 '{os.path.basename(filepath)}'에 대한 설명]:\n{desc}"]

# --- 메인 처리 함수 ---
def process_document(filepath: str) -> list[str]:
    """파일 확장자에 따라 적절한 처리 함수를 호출하는 디스패처 함수."""
    ext = os.path.splitext(filepath)[1].lower()
    
    print(f"'{filepath}' 처리 시작 (확장자: {ext})")
    
    if ext == '.pdf':
        return _process_pdf(filepath)
    elif ext == '.docx':
        return _process_docx(filepath)
    elif ext == '.pptx':
        return _process_pptx(filepath)
    elif ext == '.xlsx':
        return _process_xlsx(filepath)
    elif ext in ['.jpg', '.jpeg', '.png']:
        return _process_image(filepath)
    else:
        print(f"지원하지 않는 파일 형식입니다: {ext}")
        return []
