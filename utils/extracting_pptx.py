# pip install python-pptx

import os
from pptx import Presentation
from io import BytesIO, StringIO

def pptx_to_markdown_string(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    output = StringIO()

    output.write("# PPT 자동 변환\n\n")

    for i, slide in enumerate(prs.slides):
        output.write(f"## 슬라이드 {i + 1}\n\n")

        for shape in slide.shapes:
            # 텍스트 추출
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    output.write(f"{text}\n\n")

            # 이미지 설명
            if shape.shape_type == 13:  # picture
                output.write("[실제사진이미지(그래프,텍스트아님)]\n\n")

        output.write("---\n\n")

    return output.getvalue()


if __name__ == "__main__":
    FILE_PATH = "sample.pptx"  
    # OUTPUT_PATH = "output_from_ppt.md"  # 저장할 마크다운 파일명

    result = pptx_to_markdown_string(FILE_PATH)

    # 출력
    print(result)

    # 저장
    # with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
    #     f.write(result)

    # print(f"[변환 완료] {OUTPUT_PATH}")
