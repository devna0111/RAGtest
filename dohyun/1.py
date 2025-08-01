import os
import json
import fitz               # PDF parsing
import docx
import pandas as pd
from pptx import Presentation
from typing import List, Dict

# --------------------------
# 1. 문서 파싱 (텍스트/표/이미지)
# --------------------------

def extract_text_from_pdf(file_path: str) -> List[Dict]:
    """PDF에서 텍스트 추출"""
    doc = fitz.open(file_path)
    chunks = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            chunks.append({
                "type": "text",
                "content": text.strip(),
                "metadata": {"page": page_num+1, "position": 1}
            })
    return chunks

def extract_text_from_docx(file_path: str) -> List[Dict]:
    """DOCX에서 텍스트 추출"""
    doc = docx.Document(file_path)
    chunks = []
    for idx, para in enumerate(doc.paragraphs):
        if para.text.strip():
            chunks.append({
                "type": "text",
                "content": para.text.strip(),
                "metadata": {"page": 1, "position": idx+1}
            })
    return chunks

def extract_tables_from_xlsx(file_path: str) -> List[Dict]:
    """엑셀에서 표 추출"""
    xls = pd.ExcelFile(file_path)
    chunks = []
    for sheet in xls.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet)
        chunks.append({
            "type": "table",
            "content": df.to_json(orient='records', force_ascii=False),
            "metadata": {"sheet": sheet}
        })
    return chunks

def extract_text_from_ppt(file_path: str) -> List[Dict]:
    """PPT에서 텍스트 추출"""
    prs = Presentation(file_path)
    chunks = []
    for slide_num, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
        if texts:
            chunks.append({
                "type": "text",
                "content": " ".join(texts),
                "metadata": {"page": slide_num+1, "position": 1}
            })
    return chunks

# --------------------------
# 2. 청크 분할 (토큰 기준 or 문단 기준)
# --------------------------

def chunk_text(content: str, max_length=500) -> List[str]:
    """단순 청크 분할"""
    words = content.split()
    chunks = []
    for i in range(0, len(words), max_length):
        chunks.append(" ".join(words[i:i+max_length]))
    return chunks

# --------------------------
# 3. 문서 → JSON 스키마 통합
# --------------------------

def process_document(file_path: str) -> List[Dict]:
    ext = os.path.splitext(file_path)[1].lower()
    results = []

    if ext == '.pdf':
        results.extend(extract_text_from_pdf(file_path))
    elif ext == '.docx':
        results.extend(extract_text_from_docx(file_path))
    elif ext == '.xlsx':
        results.extend(extract_tables_from_xlsx(file_path))
    elif ext == '.pptx':
        results.extend(extract_text_from_ppt(file_path))
    else:
        print("지원하지 않는 파일 형식")

    # 청크 분할
    final_chunks = []
    for item in results:
        if item["type"] == "text":
            for idx, c in enumerate(chunk_text(item["content"], 100)):
                final_chunks.append({
                    "chunk_id": f"{item['metadata']['page']}_{item['metadata']['position']}_{idx}",
                    "type": item["type"],
                    "content": c,
                    "metadata": item["metadata"]
                })
        else:
            final_chunks.append(item)
    
    return final_chunks

# --------------------------
# 4. 보고서 작성 (템플릿 활용)
# --------------------------

def generate_report(chunks: List[Dict]) -> str:
    """단순 보고서 초안 생성"""
    report = "# 보고서 초안\n\n"
    report += "## 서론\n문서를 기반으로 작성된 요약 보고서입니다.\n\n"
    report += "## 본문\n"
    
    for c in chunks:
        if c["type"] == "text":
            report += f"- {c['content']}\n"
        elif c["type"] == "table":
            report += f"\n[표 데이터]: {c['content'][:200]}...\n"
    
    report += "\n## 결론\n데이터 분석을 통해 결론을 도출할 수 있습니다.\n"
    return report


# --------------------------
# 실행 예시
# --------------------------

if __name__ == "__main__":
    file = "sample.docx"   # 테스트 문서
    chunks = process_document(file)
    report = generate_report(chunks)
    print(report)
