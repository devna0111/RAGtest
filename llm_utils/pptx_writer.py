from pptx import Presentation
import os

def save_outline_to_pptx(outline_text: str, output_path: str = "output/presentation.pptx") -> str:
    """
    슬라이드 개요 텍스트를 기반으로 PPT 파일 생성
    - 각 슬라이드는 [슬라이드 N], 제목, 핵심 포인트 형식으로 구성되어야 함
    """
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # 제목 + 내용 슬라이드

    slides = outline_text.split("[슬라이드 ")[1:]  # 슬라이드 N 분할
    for slide_block in slides:
        lines = slide_block.strip().splitlines()
        title = lines[0].replace("제목:", "").strip()
        bullet_points = [line.strip("- ").strip() for line in lines if line.startswith("-")]

        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        content = slide.placeholders[1].text_frame
        content.clear()  # 첫 문단 제거
        for point in bullet_points:
            content.add_paragraph().text = point

    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    sample_outline = """
[슬라이드 1]
제목: 사내 업무 자동화 소개
핵심 포인트:
- 업무 효율화 필요성
- RAG 기반 자동화

[슬라이드 2]
제목: 주요 기능
핵심 포인트:
- 문서 요약 및 퀴즈 생성
- 발표 자료 자동 작성
- 벡터DB 기반 검색

[슬라이드 3]
제목: 기대 효과
핵심 포인트:
- 교육자료 제작 시간 단축
- 구성원 이해도 향상
"""
    output_path = save_outline_to_pptx(sample_outline, "output/sample_presentation.pptx")
    print(f"PPTX 저장 완료: {output_path}")
